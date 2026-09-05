import os
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".txt",
    ".md",
    ".pptx",
    ".csv",
    ".xlsx",
    ".xls",
    ".json",
}


def read_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    file_text = ""

    if ext == ".pdf":
        reader = PdfReader(file_path)
        file_text = "".join(page.extract_text() + "\n" for page in reader.pages)

    elif ext == ".docx":
        doc = Document(file_path)
        file_text = "\n".join([para.text for para in doc.paragraphs])

    elif ext in [".txt", ".md"]:
        with open(file_path, "r", encoding="utf-8") as f:
            file_text = f.read()
    elif ext == ".pptx":
        prs = Presentation(file_path)
        file_text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif ext == ".csv":
        df = pd.read_csv(file_path)
        file_text = df.to_string()

    elif ext in [".xlsx", ".xls"]:
        df = pd.read_excel(file_path)
        file_text = df.to_string()

    elif ext == ".json":
        import json
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        file_text = str(data)

    else:
        file_text = f"Unsupported file type: {ext or 'unknown'}"

    return file_text
